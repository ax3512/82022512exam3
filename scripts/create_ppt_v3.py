# -*- coding: utf-8 -*-
"""IA 설계문서 AI 분석 시스템 — AIOps 전략 발표 PPT (KTDS 스타일)."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── 컬러 팔레트 (KTDS 코퍼레이트) ──────────────────────────────
C_KT_RED    = RGBColor(0xE0, 0x00, 0x24)   # KT 레드
C_DARK      = RGBColor(0x1A, 0x1A, 0x2E)   # 다크 네이비
C_PRIMARY   = RGBColor(0x00, 0x5A, 0xB5)   # 딥 블루
C_ACCENT    = RGBColor(0x00, 0xC4, 0xB4)   # 민트
C_PURPLE    = RGBColor(0x6C, 0x5C, 0xE7)   # 퍼플
C_ORANGE    = RGBColor(0xFF, 0x6B, 0x35)
C_WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
C_LIGHT     = RGBColor(0xE8, 0xEC, 0xF1)
C_GRAY      = RGBColor(0x90, 0xA4, 0xAE)
C_DARK_GRAY = RGBColor(0x4A, 0x4A, 0x5A)
C_GREEN     = RGBColor(0x00, 0xB8, 0x94)
C_RED       = RGBColor(0xE5, 0x39, 0x35)
C_BG        = RGBColor(0xF7, 0xF8, 0xFC)
C_CARD      = RGBColor(0xFF, 0xFF, 0xFF)
C_NAVY      = RGBColor(0x0D, 0x1B, 0x3E)
C_GOLD      = RGBColor(0xFF, 0xB3, 0x00)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
W = prs.slide_width
H = prs.slide_height

TOTAL = 10

# ── 유틸 ────────────────────────────────────────────────────────
def bg(slide, color):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color

def rect(slide, l, t, w, h, fill, line=None, radius=0.05):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = fill
    if line: s.line.color.rgb = line; s.line.width = Pt(1)
    else: s.line.fill.background()
    s.adjustments[0] = radius
    return s

def box(slide, l, t, w, h, fill, text="", sz=12, color=C_WHITE, bold=False, align=PP_ALIGN.CENTER, line=None):
    s = rect(slide, l, t, w, h, fill, line)
    tf = s.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = text; p.font.size = Pt(sz)
    p.font.color.rgb = color; p.font.bold = bold; p.alignment = align
    return s

def txt(slide, l, t, w, h, text, sz=14, color=C_DARK, bold=False, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = text; p.font.size = Pt(sz)
    p.font.color.rgb = color; p.font.bold = bold; p.alignment = align
    return tb

def multi(slide, l, t, w, h, lines, sz=13, color=C_DARK, spacing=1.3, bold_first=False):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line; p.font.size = Pt(sz); p.font.color.rgb = color
        p.space_after = Pt(3); p.line_spacing = Pt(sz * spacing)
        if bold_first and i == 0: p.font.bold = True
    return tb

def title_bar(slide, title, sub=""):
    rect(slide, Inches(0), Inches(0), W, Inches(1.05), C_NAVY)
    # 좌측 레드 악센트
    rect(slide, Inches(0), Inches(0), Inches(0.06), Inches(1.05), C_KT_RED, radius=0)
    txt(slide, Inches(0.5), Inches(0.15), Inches(9), Inches(0.5), title, sz=26, color=C_WHITE, bold=True)
    if sub:
        txt(slide, Inches(0.5), Inches(0.58), Inches(10), Inches(0.4), sub, sz=13, color=C_GRAY)

def page_num(slide, n):
    txt(slide, Inches(12.2), Inches(7.05), Inches(1), Inches(0.4),
        f"{n} / {TOTAL}", sz=10, color=C_GRAY, align=PP_ALIGN.RIGHT)

def circle_icon(slide, x, y, size, color, text, sz=28):
    c = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, size, size)
    c.fill.solid(); c.fill.fore_color.rgb = color; c.line.fill.background()
    tf = c.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = text; p.font.size = Pt(sz); p.alignment = PP_ALIGN.CENTER
    return c

# ═══════════════════════════════════════════════════════════════
# SLIDE 1: 표지
# ═══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
bg(sl, C_NAVY)

# 좌측 레드라인
rect(sl, Inches(0), Inches(0), Inches(0.07), H, C_KT_RED, radius=0)
rect(sl, Inches(0.11), Inches(0), Inches(0.025), H, C_ACCENT, radius=0)

# 상단 태그
box(sl, Inches(1.2), Inches(1.2), Inches(2.2), Inches(0.42), C_KT_RED,
    "AIOps Initiative", sz=14, color=C_WHITE, bold=True)

# 메인 타이틀
txt(sl, Inches(1.2), Inches(2.0), Inches(11), Inches(0.7),
    "AI 기반 IA 설계문서", sz=44, color=C_WHITE, bold=True)
txt(sl, Inches(1.2), Inches(2.75), Inches(11), Inches(0.7),
    "자동 분석 및 Q&A 시스템", sz=44, color=C_ACCENT, bold=True)

# 구분선
rect(sl, Inches(1.2), Inches(3.65), Inches(4), Inches(0.04), C_KT_RED, radius=0)

# 부제
txt(sl, Inches(1.2), Inches(3.9), Inches(10), Inches(0.5),
    "RAG + PostgreSQL + pgvector 기반 설계문서 지식 자동화 플랫폼", sz=17, color=C_GRAY)

# 키워드 뱃지
badges = ["RAG Pipeline", "Vector Search", "LLM Automation", "AIOps", "Knowledge Base"]
for i, b in enumerate(badges):
    box(sl, Inches(1.2 + i * 2.2), Inches(4.8), Inches(2.0), Inches(0.42),
        RGBColor(0x1F, 0x2B, 0x4D), b, sz=12, color=C_ACCENT, line=C_ACCENT)

# 우하단
txt(sl, Inches(9), Inches(6.0), Inches(4), Inches(0.4),
    "K-BILL 무선 요금 시스템", sz=15, color=C_GRAY, align=PP_ALIGN.RIGHT)
txt(sl, Inches(9), Inches(6.4), Inches(4), Inches(0.4),
    "2025.03", sz=14, color=C_GRAY, align=PP_ALIGN.RIGHT)


# ═══════════════════════════════════════════════════════════════
# SLIDE 2: 추진 배경
# ═══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
bg(sl, C_BG)
title_bar(sl, "01  추진 배경", "AIOps 전략 방향 및 현장 Pain Point")
page_num(sl, 2)

# 좌측: AIOps 전략
box(sl, Inches(0.5), Inches(1.3), Inches(6.0), Inches(0.55), C_PRIMARY,
    "KTDS AIOps 전략 방향", sz=16, color=C_WHITE, bold=True, align=PP_ALIGN.LEFT)

aiops_lines = [
    "  AI를 활용한 IT 운영 자동화 및 생산성 혁신",
    "  반복적 수작업 → AI 자동화로 전환",
    "  지식 기반(Knowledge Base) 구축을 통한 업무 효율화",
    "  개발자 역량을 창의적 업무에 집중",
]
multi(sl, Inches(0.6), Inches(2.0), Inches(5.8), Inches(2.5), aiops_lines, sz=14, color=C_DARK)

# 우측: Pain Point
box(sl, Inches(6.8), Inches(1.3), Inches(6.0), Inches(0.55), C_RED,
    "현장 Pain Point", sz=16, color=C_WHITE, bold=True, align=PP_ALIGN.LEFT)

pain_items = [
    ("30~50p", "설계문서 평균 분량\n수십 건 문서 수작업 분석"),
    ("2~4h", "신규 투입 인력\n1건 분석 소요시간"),
    ("반복", "동일 질문 반복 발생\n테이블/프로그램 찾기"),
    ("유실", "담당자 변경 시\n지식 단절 발생"),
]
for i, (num, desc) in enumerate(pain_items):
    y = Inches(2.1 + i * 1.15)
    box(sl, Inches(7.0), y, Inches(1.2), Inches(0.9), C_NAVY, num, sz=22, color=C_KT_RED, bold=True)
    multi(sl, Inches(8.4), y + Inches(0.05), Inches(4.2), Inches(0.85),
          desc.split("\n"), sz=13, color=C_DARK)

# 하단: 화살표 결론
rect(sl, Inches(0.5), Inches(6.2), Inches(12.3), Inches(0.9), RGBColor(0xE8, 0xF0, 0xFE))
txt(sl, Inches(0.7), Inches(6.25), Inches(12), Inches(0.8),
    "AI가 설계문서를 읽고, 분석하고, 답변한다  →  개발자는 개발에만 집중",
    sz=18, color=C_PRIMARY, bold=True, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════
# SLIDE 3: 프로젝트 개요
# ═══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
bg(sl, C_BG)
title_bar(sl, "02  프로젝트 개요", "목표 · 범위 · 핵심 가치")
page_num(sl, 3)

# 목표
box(sl, Inches(0.5), Inches(1.3), Inches(12.3), Inches(0.55), C_PRIMARY,
    "프로젝트 목표:  IA 설계문서 지식을 AI가 자동으로 구조화하고, 자연어로 즉시 검색·답변하는 시스템 구축",
    sz=15, color=C_WHITE, bold=True, align=PP_ALIGN.LEFT)

# 4개 핵심 기능 카드
features = [
    ("01", "자동 파싱", "DOCX 설계문서를\n섹션/테이블 단위로\n자동 구조화 파싱", C_PRIMARY),
    ("02", "AI 요약", "LLM이 각 섹션의\n핵심 내용을 요약하고\n테이블명·JO명 자동 추출", C_PURPLE),
    ("03", "벡터 검색", "의미 기반 유사도 검색\npgvector 코사인 매칭\ntop-K 검색 엔진", C_GREEN),
    ("04", "Q&A 챗봇", "자연어 질문에\n출처 기반 정확한 답변\n할루시네이션 방지", C_ORANGE),
]
for i, (num, title, desc, color) in enumerate(features):
    x = Inches(0.5 + i * 3.15)
    card = rect(sl, x, Inches(2.2), Inches(2.9), Inches(3.2), C_CARD, C_LIGHT)
    # 번호 서클
    circle_icon(sl, x + Inches(0.95), Inches(2.4), Inches(0.9), color, num, sz=24)
    # 제목
    txt(sl, x + Inches(0.1), Inches(3.45), Inches(2.7), Inches(0.4),
        title, sz=17, color=color, bold=True, align=PP_ALIGN.CENTER)
    # 설명
    multi(sl, x + Inches(0.2), Inches(3.9), Inches(2.5), Inches(1.5),
          desc.split("\n"), sz=12, color=C_DARK_GRAY, spacing=1.4)

# 하단: 적용 범위
box(sl, Inches(0.5), Inches(5.7), Inches(3.8), Inches(1.2), RGBColor(0xF0, 0xF4, 0xF8),
    "", line=C_LIGHT)
txt(sl, Inches(0.7), Inches(5.8), Inches(3.5), Inches(0.3),
    "적용 범위", sz=13, color=C_PRIMARY, bold=True)
multi(sl, Inches(0.7), Inches(6.15), Inches(3.5), Inches(0.7),
      ["K-BILL 무선 요금 시스템", "IA 설계문서 (.docx)"], sz=12, color=C_DARK)

box(sl, Inches(4.6), Inches(5.7), Inches(3.8), Inches(1.2), RGBColor(0xF0, 0xF4, 0xF8),
    "", line=C_LIGHT)
txt(sl, Inches(4.8), Inches(5.8), Inches(3.5), Inches(0.3),
    "기술 스택", sz=13, color=C_PRIMARY, bold=True)
multi(sl, Inches(4.8), Inches(6.15), Inches(3.5), Inches(0.7),
      ["PostgreSQL 16 + pgvector", "GPT-4o-mini / E5 임베딩"], sz=12, color=C_DARK)

box(sl, Inches(8.7), Inches(5.7), Inches(4.1), Inches(1.2), RGBColor(0xF0, 0xF4, 0xF8),
    "", line=C_LIGHT)
txt(sl, Inches(8.9), Inches(5.8), Inches(3.8), Inches(0.3),
    "핵심 가치", sz=13, color=C_PRIMARY, bold=True)
multi(sl, Inches(8.9), Inches(6.15), Inches(3.8), Inches(0.7),
      ["수작업 분석 → AI 자동화", "지식 유실 → 영구 보존"], sz=12, color=C_DARK)


# ═══════════════════════════════════════════════════════════════
# SLIDE 4: 시스템 아키텍처
# ═══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
bg(sl, C_BG)
title_bar(sl, "03  시스템 아키텍처", "RAG 기반 End-to-End 파이프라인")
page_num(sl, 4)

# 좌측: 적재 파이프라인
txt(sl, Inches(0.5), Inches(1.3), Inches(5), Inches(0.4),
    "Document Ingestion Pipeline", sz=16, color=C_PRIMARY, bold=True)

load_steps = [
    ("DOCX", "설계문서\n업로드", RGBColor(0x42, 0x6F, 0xB1)),
    ("Parser", "섹션 분할\n메타 추출", RGBColor(0x5C, 0x6B, 0xC0)),
    ("LLM", "AI 요약\n엔티티 추출", C_PURPLE),
    ("Embedder", "벡터 변환\n1024차원", C_GREEN),
    ("PG+pgvec", "통합 저장\nHNSW 인덱스", C_PRIMARY),
]
for i, (label, desc, color) in enumerate(load_steps):
    x = Inches(0.5 + i * 2.45)
    y = Inches(1.85)
    box(sl, x, y, Inches(2.1), Inches(0.55), color, label, sz=14, color=C_WHITE, bold=True)
    multi(sl, x + Inches(0.1), y + Inches(0.65), Inches(1.9), Inches(0.7),
          desc.split("\n"), sz=11, color=C_DARK_GRAY)
    if i < len(load_steps) - 1:
        txt(sl, x + Inches(2.0), y + Inches(0.1), Inches(0.5), Inches(0.4),
            "→", sz=22, color=C_GRAY, align=PP_ALIGN.CENTER)

# 구분선
rect(sl, Inches(0.5), Inches(3.35), Inches(12.3), Inches(0.03), C_LIGHT, radius=0)

# 질의 파이프라인 — 실제 호출 흐름
txt(sl, Inches(0.5), Inches(3.5), Inches(8), Inches(0.4),
    "Query & Answer Pipeline (실제 호출 흐름)", sz=16, color=C_GREEN, bold=True)

# 6단계 실제 흐름
q_steps = [
    ("1. 질문 입력", "자연어 질문", C_NAVY),
    ("2. 임베딩", "질문 → 1024차원\n벡터 변환 (e5)", C_PURPLE),
    ("3. pgvector", "코사인 유사도\ntop50 검색", C_GREEN),
    ("4. DR 추출", "DR별 best score\n상위 5개 DR\n전체 섹션 조회", RGBColor(0x00, 0x96, 0x88)),
    ("5. 1차 LLM", "DR별 섹션 필터링\n관련 섹션만 추출\n(DR당 1회 호출)", C_ORANGE),
    ("6. 2차 LLM", "필터링 결과 합산\n최종 답변 생성\n(출처 포함)", C_KT_RED),
]
for i, (label, desc, color) in enumerate(q_steps):
    x = Inches(0.3 + i * 2.15)
    y = Inches(3.95)
    box(sl, x, y, Inches(1.9), Inches(0.45), color, label, sz=11, color=C_WHITE, bold=True)
    multi(sl, x + Inches(0.05), y + Inches(0.5), Inches(1.8), Inches(1.1),
          desc.split('\n'), sz=10, color=C_DARK_GRAY)
    if i < len(q_steps) - 1:
        txt(sl, x + Inches(1.8), y + Inches(0.05), Inches(0.4), Inches(0.35),
            '→', sz=16, color=C_GRAY, align=PP_ALIGN.CENTER)

# 핵심 포인트 박스
rect(sl, Inches(0.5), Inches(5.5), Inches(12.3), Inches(0.55), RGBColor(0xFE, 0xF3, 0xE2))
txt(sl, Inches(0.7), Inches(5.55), Inches(12), Inches(0.45),
    '핵심: 질문 임베딩 → pgvector top50 → 상위 5개 DR 전체 섹션 → 1차 LLM DR별 필터링 → 2차 LLM 최종 답변 (2단계 LLM)',
    sz=13, color=C_ORANGE, bold=True)

# 하단: DB 구조
txt(sl, Inches(0.5), Inches(6.15), Inches(5), Inches(0.35),
    'Unified Database (PostgreSQL 16 + pgvector)', sz=14, color=C_NAVY, bold=True)

db_tables = [
    ('documents', '문서 메타 | DR번호, 제목, 시스템, 종합요약', C_PRIMARY),
    ('sections', '섹션 + 벡터 | 내용, AI요약, embedding(1024), 엔티티', C_GREEN),
    ('feedback', '피드백 | 답변 평가, 개선 데이터', C_PURPLE),
]
for i, (name, desc, color) in enumerate(db_tables):
    x = Inches(0.5 + i * 4.2)
    box(sl, x, Inches(6.55), Inches(3.8), Inches(0.65), C_CARD, f'{name}:  {desc}',
        sz=11, color=color, bold=True, align=PP_ALIGN.LEFT, line=color)


# ═══════════════════════════════════════════════════════════════
# SLIDE 5: 핵심 AI 기술
# ═══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
bg(sl, C_BG)
title_bar(sl, "04  핵심 AI 기술", "RAG (Retrieval-Augmented Generation) 아키텍처")
page_num(sl, 5)

# RAG 설명
box(sl, Inches(0.5), Inches(1.3), Inches(12.3), Inches(0.9), RGBColor(0xEA, 0xF0, 0xFF),
    "", line=C_PRIMARY)
multi(sl, Inches(0.7), Inches(1.35), Inches(12), Inches(0.8),
      ["RAG (Retrieval-Augmented Generation): LLM이 답변할 때 외부 지식(설계문서)을 실시간으로 검색하여 참조하는 기술",
       "→ 할루시네이션 방지 + 출처 명시 + 최신 문서 반영 가능"],
      sz=14, color=C_PRIMARY)

# 3개 기술 블록
techs = [
    ("벡터 임베딩", "multilingual-e5-large",
     [
         "문장의 의미를 1024차원 벡터로 변환",
         "다국어 지원 (한국어 최적화)",
         "설계문서 전문 용어 의미 포착",
         "문서 적재 시 1회 변환 → DB 저장",
     ], C_PURPLE),
    ("pgvector 검색", "HNSW 인덱스 + Cosine",
     [
         "PostgreSQL 내장 벡터 검색",
         "HNSW: 근사 최근접 이웃 알고리즘",
         "Cosine 유사도 기반 의미 매칭",
         "별도 벡터DB 불필요 → 운영 단순화",
     ], C_GREEN),
    ("LLM 답변 생성", "GPT-4o-mini",
     [
         "검색 DR의 전체 섹션 1차 필터링 후 답변",
         "출처(DR번호, 섹션 경로) 명시",
         "테이블명/JO명 원문 그대로 인용",
         "2단계 필터링으로 정확도 향상",
     ], C_ORANGE),
]

for i, (title, subtitle, items, color) in enumerate(techs):
    x = Inches(0.5 + i * 4.2)
    y = Inches(2.6)
    box(sl, x, y, Inches(3.8), Inches(0.5), color, title, sz=16, color=C_WHITE, bold=True)
    txt(sl, x, y + Inches(0.55), Inches(3.8), Inches(0.3),
        subtitle, sz=11, color=C_GRAY, align=PP_ALIGN.CENTER)
    card = rect(sl, x, y + Inches(0.9), Inches(3.8), Inches(2.8), C_CARD, C_LIGHT)
    bullet_lines = [f"  {item}" for item in items]
    multi(sl, x + Inches(0.15), y + Inches(1.0), Inches(3.5), Inches(2.6),
          bullet_lines, sz=12, color=C_DARK)

# 하단 차별점
box(sl, Inches(0.5), Inches(6.5), Inches(12.3), Inches(0.7), C_NAVY,
    "차별점:  단일 PostgreSQL에 문서·벡터·메타데이터 통합  |  별도 벡터DB/인프라 불필요  |  ACID 트랜잭션 보장  |  운영 복잡도 최소화",
    sz=14, color=C_ACCENT, bold=True, align=PP_ALIGN.LEFT)


# ═══════════════════════════════════════════════════════════════
# SLIDE 6: 주요 기능 상세
# ═══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
bg(sl, C_BG)
title_bar(sl, "05  주요 기능 상세", "설계문서 분석 → 지식 검색 → 업무 활용")
page_num(sl, 6)

funcs = [
    ("자동 구조화 파싱", [
        "DOCX 파일 자동 분석",
        "제목 계층 → 섹션 트리 구성",
        "메타 정보 자동 추출 (DR번호, 시스템, 대상년월)",
        "마크다운 테이블 변환",
    ], C_PRIMARY),
    ("AI 요약 + 엔티티 추출", [
        "섹션 유형별 맞춤 프롬프트 (과제분석/구현방안/검증방안/DB참조)",
        "테이블명(XX_YYY_ZZZ), JO명(XxxYyyJO) 자동 추출",
        "mentioned_tables / mentioned_sources → JSONB 저장",
        "빈 컨텐츠/빈 표 감지 시 LLM 미호출 (할루시네이션 방지)",
    ], C_PURPLE),
    ("벡터 유사도 검색", [
        "pgvector 코사인 유사도 검색",
        "DR번호 지정 시 범위 필터링",
        "DR별 best score 상위 5개 추출",
        "해당 DR 전체 섹션 자동 조회",
    ], C_GREEN),
    ("출처 기반 Q&A", [
        "자연어 질문 → AI 답변",
        "DR번호 + 섹션 경로 출처 명시",
        "관련 없는 문서 자동 필터링",
        "피드백 수집 → DB 저장",
    ], C_ORANGE),
    ("문서 관리", [
        "문서 목록 조회 / 검색",
        "개별 삭제 (CASCADE 자동 정리)",
        "폴더 일괄 적재 (스트리밍 진행률)",
        "섹션 요약 수동 편집 → 벡터 재생성",
    ], C_NAVY),
]

for i, (title, items, color) in enumerate(funcs):
    x = Inches(0.5 + (i % 3) * 4.2)
    y = Inches(1.3) if i < 3 else Inches(4.4)
    box(sl, x, y, Inches(3.8), Inches(0.5), color, title, sz=14, color=C_WHITE, bold=True, align=PP_ALIGN.LEFT)
    card = rect(sl, x, y + Inches(0.55), Inches(3.8), Inches(2.3) if i < 3 else Inches(2.3), C_CARD, C_LIGHT)
    bullet_lines = [f"  {item}" for item in items]
    multi(sl, x + Inches(0.15), y + Inches(0.65), Inches(3.5), Inches(2.1),
          bullet_lines, sz=12, color=C_DARK)


# ═══════════════════════════════════════════════════════════════
# SLIDE 7: 기대효과 — 정량적
# ═══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
bg(sl, C_BG)
title_bar(sl, "06  기대효과 — 생산성 향상", "정량적 개선 지표")
page_num(sl, 7)

# 큰 숫자 카드
metrics = [
    ("문서 분석 시간", "2~4시간", "5분 이내", "95%↓", C_PRIMARY),
    ("신규 투입 인력\n업무 파악", "2~3일", "즉시 질의", "90%↓", C_GREEN),
    ("테이블/JO\n검색 시간", "30분~1시간", "3초", "99%↓", C_PURPLE),
    ("지식 유실률", "담당자 퇴사 시\n100% 유실", "0%\n(DB 영구 보존)", "100%↓", C_KT_RED),
]

for i, (title, before, after, pct, color) in enumerate(metrics):
    x = Inches(0.5 + i * 3.15)
    # 카드
    card = rect(sl, x, Inches(1.3), Inches(2.9), Inches(4.8), C_CARD, C_LIGHT)
    # 제목
    multi(sl, x + Inches(0.15), Inches(1.4), Inches(2.6), Inches(0.7),
          title.split("\n"), sz=14, color=C_DARK, bold_first=True)
    # 개선율 큰 숫자
    txt(sl, x + Inches(0.1), Inches(2.2), Inches(2.7), Inches(0.8),
        pct, sz=44, color=color, bold=True, align=PP_ALIGN.CENTER)
    # Before
    box(sl, x + Inches(0.15), Inches(3.2), Inches(2.6), Inches(0.35), RGBColor(0xFF, 0xEB, 0xEE),
        "AS-IS", sz=11, color=C_RED, bold=True)
    multi(sl, x + Inches(0.15), Inches(3.6), Inches(2.6), Inches(0.7),
          before.split("\n"), sz=12, color=C_RED)
    # After
    box(sl, x + Inches(0.15), Inches(4.4), Inches(2.6), Inches(0.35), RGBColor(0xE8, 0xF5, 0xE9),
        "TO-BE", sz=11, color=C_GREEN, bold=True)
    multi(sl, x + Inches(0.15), Inches(4.8), Inches(2.6), Inches(0.7),
          after.split("\n"), sz=12, color=C_GREEN)

# 하단 ROI
rect(sl, Inches(0.5), Inches(6.35), Inches(12.3), Inches(0.85), C_NAVY)
txt(sl, Inches(0.7), Inches(6.4), Inches(12), Inches(0.4),
    "연간 예상 절감 효과", sz=14, color=C_GRAY)
txt(sl, Inches(0.7), Inches(6.75), Inches(12), Inches(0.4),
    "설계문서 분석 공수 연 960시간 절감 (20건/월 × 4h × 12월)  →  인력 0.5명 효과",
    sz=16, color=C_ACCENT, bold=True)


# ═══════════════════════════════════════════════════════════════
# SLIDE 8: 기대효과 — AIOps 성숙도
# ═══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
bg(sl, C_BG)
title_bar(sl, "07  AIOps 성숙도 기여", "AI 자동화 단계별 로드맵")
page_num(sl, 8)

# 성숙도 5단계 바
stages = [
    ("Level 1", "수동", "사람이 직접\n문서 분석", C_GRAY, False),
    ("Level 2", "보조", "키워드 검색\n단순 필터링", RGBColor(0x90, 0xCA, 0xF9), False),
    ("Level 3", "자동화", "AI 요약/검색\n자연어 Q&A", C_GREEN, True),  # 현재
    ("Level 4", "예측", "변경영향도 예측\n자동 TC 생성", C_PRIMARY, False),
    ("Level 5", "자율", "자동 설계 검증\n이상 탐지", C_PURPLE, False),
]

for i, (level, name, desc, color, current) in enumerate(stages):
    x = Inches(0.5 + i * 2.5)
    y = Inches(1.5)
    h = Inches(2.8) if current else Inches(2.4)
    # 카드
    c = rect(sl, x, y, Inches(2.2), h, color if current else C_CARD,
             color if not current else None)
    if current:
        # 현재 위치 뱃지
        box(sl, x + Inches(0.3), y - Inches(0.3), Inches(1.6), Inches(0.35), C_KT_RED,
            "현재 단계", sz=12, color=C_WHITE, bold=True)
    font_c = C_WHITE if current else C_DARK
    txt(sl, x + Inches(0.1), y + Inches(0.15), Inches(2.0), Inches(0.3),
        level, sz=12, color=font_c if current else C_GRAY, bold=True, align=PP_ALIGN.CENTER)
    txt(sl, x + Inches(0.1), y + Inches(0.45), Inches(2.0), Inches(0.35),
        name, sz=20, color=font_c if current else color, bold=True, align=PP_ALIGN.CENTER)
    multi(sl, x + Inches(0.15), y + Inches(0.95), Inches(1.9), Inches(1.0),
          desc.split("\n"), sz=12, color=font_c if current else C_DARK_GRAY)

# 하단: 본 프로젝트 기여 영역
txt(sl, Inches(0.5), Inches(4.7), Inches(12), Inches(0.4),
    "본 프로젝트의 AIOps 기여 영역", sz=18, color=C_NAVY, bold=True)

contribs = [
    ("지식 자산화", "개인 지식 → 조직 지식\n설계문서 DB 영구 보존\n담당자 변경에도 지식 유지", C_PRIMARY),
    ("업무 자동화", "문서 분석 수작업 제거\nAI 요약으로 분석 시간 95%↓\n반복 질의 자동 응답", C_GREEN),
    ("의사결정 지원", "테이블 영향도 즉시 파악\n관련 문서 자동 연결\n변경 리스크 사전 인지", C_PURPLE),
    ("확장 기반 구축", "PostgreSQL 표준 인프라\nRAG 파이프라인 재사용 가능\n타 시스템 설계문서 확장", C_ORANGE),
]

for i, (title, desc, color) in enumerate(contribs):
    x = Inches(0.5 + i * 3.15)
    y = Inches(5.2)
    box(sl, x, y, Inches(2.9), Inches(0.45), color, title, sz=14, color=C_WHITE, bold=True)
    card = rect(sl, x, y + Inches(0.5), Inches(2.9), Inches(1.5), C_CARD, C_LIGHT)
    multi(sl, x + Inches(0.15), y + Inches(0.6), Inches(2.6), Inches(1.3),
          desc.split("\n"), sz=12, color=C_DARK)


# ═══════════════════════════════════════════════════════════════
# SLIDE 9: 확대 적용 계획
# ═══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
bg(sl, C_BG)
title_bar(sl, "08  확대 적용 계획", "단계별 확산 로드맵")
page_num(sl, 9)

phases = [
    ("Phase 1", "현재", "K-BILL 무선 적용", [
        "IA 설계문서 적재 완료",
        "Q&A 챗봇 운영 시작",
        "사용자 피드백 수집",
        "검색 정확도 튜닝",
    ], C_GREEN),
    ("Phase 2", "1~3개월", "팀 내 확산", [
        "전체 설계문서 일괄 적재",
        "Oracle DB 연동 강화",
        "하이브리드 검색 고도화",
        "답변 품질 개선",
    ], C_PRIMARY),
    ("Phase 3", "3~6개월", "본부 확대", [
        "유선/IPTV 등 타 시스템 확장",
        "변경영향도 분석 기능",
        "자동 테스트케이스 생성",
        "Slack/Teams 연동",
    ], C_PURPLE),
    ("Phase 4", "6개월~", "전사 플랫폼화", [
        "표준 RAG 플랫폼으로 고도화",
        "타 사업부 적용",
        "MCP 연동 (IDE 통합)",
        "AIOps Level 4 진입",
    ], C_ORANGE),
]

for i, (phase, period, title, items, color) in enumerate(phases):
    x = Inches(0.5 + i * 3.15)
    # 헤더
    box(sl, x, Inches(1.3), Inches(2.9), Inches(0.5), color,
        f"{phase}  |  {period}", sz=14, color=C_WHITE, bold=True)
    txt(sl, x, Inches(1.9), Inches(2.9), Inches(0.35),
        title, sz=15, color=color, bold=True, align=PP_ALIGN.CENTER)
    # 항목 카드
    card = rect(sl, x, Inches(2.35), Inches(2.9), Inches(2.8), C_CARD, C_LIGHT)
    bullet_lines = [f"  {item}" for item in items]
    multi(sl, x + Inches(0.15), Inches(2.5), Inches(2.6), Inches(2.6),
          bullet_lines, sz=12, color=C_DARK)

# 화살표 연결
for i in range(3):
    x = Inches(3.35 + i * 3.15)
    txt(sl, x, Inches(2.9), Inches(0.5), Inches(0.5),
        "→", sz=28, color=C_GRAY, align=PP_ALIGN.CENTER)

# 하단 비전
rect(sl, Inches(0.5), Inches(5.5), Inches(12.3), Inches(1.6), C_NAVY)
txt(sl, Inches(0.7), Inches(5.6), Inches(12), Inches(0.4),
    "Vision", sz=14, color=C_GOLD, bold=True)
txt(sl, Inches(0.7), Inches(6.0), Inches(11.5), Inches(0.5),
    "\"설계문서를 읽는 시간을 제로로 만들고, 모든 개발 지식을 AI가 관리하는 AIOps 환경 구현\"",
    sz=20, color=C_WHITE, bold=True)
txt(sl, Inches(0.7), Inches(6.55), Inches(11.5), Inches(0.4),
    "단일 시스템 도구 → 본부 표준 → 전사 AI Knowledge Platform 으로 확장",
    sz=14, color=C_ACCENT)


# ═══════════════════════════════════════════════════════════════
# SLIDE 10: 마무리
# ═══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
bg(sl, C_NAVY)

rect(sl, Inches(0), Inches(0), Inches(0.07), H, C_KT_RED, radius=0)

txt(sl, Inches(1.2), Inches(1.5), Inches(11), Inches(0.5),
    "Summary", sz=16, color=C_GOLD, bold=True)

summary_items = [
    "IA 설계문서를 AI가 자동으로 파싱·요약·검색·답변하는 시스템 구축",
    "PostgreSQL + pgvector 단일 DB로 문서·벡터·메타데이터 통합 관리",
    "설계문서 분석 시간 95% 절감, 연 960시간 공수 절감 기대",
    "조직 지식 영구 보존 → 담당자 변경에도 지식 단절 Zero",
    "AIOps Level 3 달성 → Level 4 (예측/자동 검증) 진입 기반 마련",
]
for i, item in enumerate(summary_items):
    y = Inches(2.1 + i * 0.65)
    rect(sl, Inches(1.2), y + Inches(0.08), Inches(0.08), Inches(0.25), C_ACCENT, radius=0)
    txt(sl, Inches(1.5), y, Inches(10), Inches(0.5),
        item, sz=17, color=C_WHITE)

# 하단 메시지
rect(sl, Inches(1.2), Inches(5.5), Inches(5), Inches(0.04), C_KT_RED, radius=0)
txt(sl, Inches(1.2), Inches(5.8), Inches(11), Inches(0.5),
    "AI로 일하는 방식을 바꿉니다.", sz=28, color=C_ACCENT, bold=True)

txt(sl, Inches(1.2), Inches(6.5), Inches(4), Inches(0.4),
    "감사합니다.", sz=20, color=C_GRAY)

page_num(sl, 10)


# ═══════════════════════════════════════════════════════════════
out = r"C:\Users\my\ia-chatbot-v3\IA_AI_QA_System_v3.pptx"
prs.save(out)
print(f"Done: {out}")
print(f"Slides: {TOTAL}")
