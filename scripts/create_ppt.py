# -*- coding: utf-8 -*-
"""IA 챗봇 V3 진급심사 발표 PPT 생성 스크립트."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ── 컬러 팔레트 ────────────────────────────────────────────────
C_DARK      = RGBColor(0x1B, 0x1B, 0x2F)   # 남색 배경
C_PRIMARY   = RGBColor(0x00, 0x6E, 0xD6)   # KT 블루
C_ACCENT    = RGBColor(0x00, 0xC4, 0xB4)   # 민트 강조
C_ORANGE    = RGBColor(0xFF, 0x6B, 0x35)   # 오렌지
C_WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
C_LIGHT     = RGBColor(0xE8, 0xEC, 0xF1)
C_GRAY      = RGBColor(0x90, 0xA4, 0xAE)
C_DARK_GRAY = RGBColor(0x4A, 0x4A, 0x5A)
C_GREEN     = RGBColor(0x4C, 0xAF, 0x50)
C_RED       = RGBColor(0xE5, 0x39, 0x35)
C_YELLOW    = RGBColor(0xFF, 0xC1, 0x07)
C_BG_CARD   = RGBColor(0x25, 0x25, 0x3A)   # 카드 배경
C_BG_SLIDE  = RGBColor(0xF5, 0xF7, 0xFA)   # 밝은 배경

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
W = prs.slide_width
H = prs.slide_height


# ── 유틸리티 함수 ──────────────────────────────────────────────
def add_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_rect(slide, left, top, width, height, fill_color, line_color=None, line_width=Pt(0)):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = line_width
    else:
        shape.line.fill.background()
    # 모서리 둥글기
    shape.adjustments[0] = 0.05
    return shape

def add_box(slide, left, top, width, height, fill_color, text="", font_size=12, font_color=C_WHITE, bold=False, align=PP_ALIGN.CENTER, line_color=None):
    shape = add_rect(slide, left, top, width, height, fill_color, line_color)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = font_color
    p.font.bold = bold
    p.alignment = align
    tf.paragraphs[0].space_before = Pt(0)
    tf.paragraphs[0].space_after = Pt(0)
    return shape

def add_text(slide, left, top, width, height, text, font_size=14, font_color=C_DARK, bold=False, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = font_color
    p.font.bold = bold
    p.alignment = align
    return txBox

def add_multiline(slide, left, top, width, height, lines, font_size=13, font_color=C_DARK, bold_first=False, line_spacing=1.3):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(font_size)
        p.font.color.rgb = font_color
        p.space_after = Pt(4)
        if bold_first and i == 0:
            p.font.bold = True
        p.line_spacing = Pt(font_size * line_spacing)
    return txBox

def add_arrow(slide, x1, y1, x2, y2, color=C_PRIMARY, width=Pt(3)):
    connector = slide.shapes.add_connector(1, x1, y1, x2, y2)  # straight
    connector.line.color.rgb = color
    connector.line.width = width
    return connector

def add_line_shape(slide, x1, y1, x2, y2, color=C_LIGHT, width=Pt(1)):
    connector = slide.shapes.add_connector(1, x1, y1, x2, y2)
    connector.line.color.rgb = color
    connector.line.width = width
    return connector

def slide_title_bar(slide, title, subtitle=""):
    """상단 타이틀 바"""
    add_rect(slide, Inches(0), Inches(0), W, Inches(1.05), C_PRIMARY)
    add_text(slide, Inches(0.6), Inches(0.15), Inches(8), Inches(0.5), title, 
             font_size=26, font_color=C_WHITE, bold=True)
    if subtitle:
        add_text(slide, Inches(0.6), Inches(0.60), Inches(10), Inches(0.4), subtitle,
                 font_size=13, font_color=RGBColor(0xBB, 0xDE, 0xFB))

def page_number(slide, num, total):
    add_text(slide, Inches(12.2), Inches(7.05), Inches(1), Inches(0.4), 
             f"{num} / {total}", font_size=10, font_color=C_GRAY, align=PP_ALIGN.RIGHT)


TOTAL_SLIDES = 8


# ═══════════════════════════════════════════════════════════════
# SLIDE 1: 표지
# ═══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_bg(sl, C_DARK)

# 좌측 장식 라인
add_rect(sl, Inches(0), Inches(0), Inches(0.08), H, C_PRIMARY)
add_rect(sl, Inches(0.12), Inches(0), Inches(0.03), H, C_ACCENT)

# 메인 타이틀
add_text(sl, Inches(1.2), Inches(1.8), Inches(11), Inches(0.6),
         "IA 설계문서 Q&A 챗봇 시스템", font_size=42, font_color=C_WHITE, bold=True)
add_text(sl, Inches(1.2), Inches(2.5), Inches(11), Inches(0.5),
         "DB 아키텍처 통합 및 고도화 (V2 → V3)", font_size=24, font_color=C_ACCENT)

# 구분선
add_rect(sl, Inches(1.2), Inches(3.3), Inches(3), Inches(0.04), C_PRIMARY)

# 하단 정보
add_text(sl, Inches(1.2), Inches(3.7), Inches(10), Inches(0.4),
         "TinyDB + ChromaDB  →  PostgreSQL 16 + pgvector 통합", font_size=16, font_color=C_GRAY)

# 키워드 뱃지
badges = ["PostgreSQL 16", "pgvector", "RAG", "FastAPI", "LLM"]
for i, badge in enumerate(badges):
    add_box(sl, Inches(1.2 + i * 2.1), Inches(4.5), Inches(1.9), Inches(0.45),
            C_BG_CARD, badge, font_size=13, font_color=C_ACCENT, line_color=C_ACCENT)

# 우하단 날짜/부서
add_text(sl, Inches(9), Inches(6.2), Inches(4), Inches(0.4),
         "2025.03", font_size=16, font_color=C_GRAY, align=PP_ALIGN.RIGHT)
add_text(sl, Inches(9), Inches(6.6), Inches(4), Inches(0.4),
         "K-BILL 무선 요금 시스템", font_size=14, font_color=C_GRAY, align=PP_ALIGN.RIGHT)


# ═══════════════════════════════════════════════════════════════
# SLIDE 2: 목차
# ═══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl, C_BG_SLIDE)
slide_title_bar(sl, "CONTENTS", "발표 목차")
page_number(sl, 2, TOTAL_SLIDES)

toc_items = [
    ("01", "프로젝트 개요",       "배경, 목적, 핵심 기능"),
    ("02", "AS-IS 아키텍처 (V2)",  "기존 시스템 구조 및 한계점"),
    ("03", "TO-BE 아키텍처 (V3)",  "PostgreSQL + pgvector 통합 설계"),
    ("04", "핵심 구현 — 호출 흐름", "질문→검색→응답 파이프라인"),
    ("05", "핵심 구현 — DB 통합",  "TinyDB+ChromaDB → PostgreSQL 마이그레이션"),
    ("06", "Before / After 비교",  "정량적·정성적 개선 효과"),
    ("07", "시연 결과",            "실제 문서 적재 및 질의 응답"),
    ("08", "향후 계획",            "확장 로드맵"),
]

for i, (num, title, desc) in enumerate(toc_items):
    y = Inches(1.4 + i * 0.72)
    # 번호 원
    add_box(sl, Inches(1.5), y, Inches(0.65), Inches(0.55), C_PRIMARY, num, 
            font_size=18, font_color=C_WHITE, bold=True)
    # 제목
    add_text(sl, Inches(2.4), y + Inches(0.02), Inches(4), Inches(0.35), title,
             font_size=18, font_color=C_DARK, bold=True)
    # 설명
    add_text(sl, Inches(2.4), y + Inches(0.32), Inches(6), Inches(0.25), desc,
             font_size=12, font_color=C_GRAY)
    # 하단 구분선
    if i < len(toc_items) - 1:
        add_line_shape(sl, Inches(1.5), y + Inches(0.65), Inches(11.5), y + Inches(0.65))


# ═══════════════════════════════════════════════════════════════
# SLIDE 3: 프로젝트 개요
# ═══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl, C_BG_SLIDE)
slide_title_bar(sl, "01  프로젝트 개요", "배경 · 목적 · 핵심 기능")
page_number(sl, 3, TOTAL_SLIDES)

# 배경 카드
add_box(sl, Inches(0.6), Inches(1.3), Inches(5.8), Inches(0.5), C_ORANGE,
        "💡 프로젝트 배경", font_size=16, font_color=C_WHITE, bold=True, align=PP_ALIGN.LEFT)
bg_lines = [
    "•  IA 설계문서는 평균 20~50페이지의 .docx 파일",
    "•  수십 건의 설계문서에서 원하는 내용 찾기 어려움",
    "•  테이블명, 프로그램ID, 비즈니스 룰을 빠르게 검색할 필요",
    "•  기존 수작업 → AI 기반 자동 분석/검색 시스템 구축",
]
add_multiline(sl, Inches(0.8), Inches(1.95), Inches(5.4), Inches(2.5), bg_lines,
              font_size=14, font_color=C_DARK)

# 목적 카드
add_box(sl, Inches(7), Inches(1.3), Inches(5.8), Inches(0.5), C_PRIMARY,
        "🎯 프로젝트 목적", font_size=16, font_color=C_WHITE, bold=True, align=PP_ALIGN.LEFT)
obj_lines = [
    "•  설계문서 자동 파싱 및 섹션별 구조화",
    "•  LLM 기반 섹션 요약 자동 생성",
    "•  벡터 유사도 검색 + 키워드 검색 결합",
    "•  자연어 질문 → 정확한 답변 + 출처 제공",
]
add_multiline(sl, Inches(7.2), Inches(1.95), Inches(5.4), Inches(2.5), obj_lines,
              font_size=14, font_color=C_DARK)

# 하단: 핵심 기능 4개 카드
features = [
    ("📄", "문서 파싱", ".docx 자동 분석\n섹션/테이블 추출"),
    ("🤖", "AI 요약", "GPT-4o-mini\n섹션별 자동 요약"),
    ("🔍", "하이브리드 검색", "벡터 유사도\n+ 키워드 매칭"),
    ("💬", "Q&A 챗봇", "자연어 질문\n출처 기반 답변"),
]
for i, (icon, title, desc) in enumerate(features):
    x = Inches(0.6 + i * 3.15)
    card = add_rect(sl, x, Inches(4.7), Inches(2.9), Inches(2.3), C_WHITE, C_LIGHT, Pt(1))
    add_text(sl, x + Inches(0.2), Inches(4.85), Inches(2.5), Inches(0.5),
             icon, font_size=32, align=PP_ALIGN.CENTER)
    add_text(sl, x + Inches(0.2), Inches(5.35), Inches(2.5), Inches(0.4),
             title, font_size=16, font_color=C_PRIMARY, bold=True, align=PP_ALIGN.CENTER)
    add_multiline(sl, x + Inches(0.3), Inches(5.8), Inches(2.3), Inches(1.0),
                  desc.split("\n"), font_size=12, font_color=C_DARK_GRAY)


# ═══════════════════════════════════════════════════════════════
# SLIDE 4: AS-IS 아키텍처 (V2)
# ═══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl, C_BG_SLIDE)
slide_title_bar(sl, "02  AS-IS 아키텍처 (V2)", "기존 시스템 구조 및 한계점")
page_number(sl, 4, TOTAL_SLIDES)

# Client 계층
add_box(sl, Inches(4.5), Inches(1.3), Inches(4.3), Inches(0.65), C_PRIMARY,
        "🖥️  Client (Web UI — app.js)", font_size=15, font_color=C_WHITE, bold=True)

# API 계층
add_box(sl, Inches(4.5), Inches(2.3), Inches(4.3), Inches(0.65), RGBColor(0x30, 0x5A, 0x9E),
        "⚡  FastAPI Server (port 8004)", font_size=15, font_color=C_WHITE, bold=True)

# Engine 계층
add_box(sl, Inches(2.5), Inches(3.3), Inches(2.6), Inches(0.55), RGBColor(0x5C, 0x6B, 0xC0),
        "search.py", font_size=13, font_color=C_WHITE)
add_box(sl, Inches(5.3), Inches(3.3), Inches(2.6), Inches(0.55), RGBColor(0x5C, 0x6B, 0xC0),
        "answer.py (LLM)", font_size=13, font_color=C_WHITE)
add_box(sl, Inches(8.1), Inches(3.3), Inches(2.6), Inches(0.55), RGBColor(0x5C, 0x6B, 0xC0),
        "summarizer.py", font_size=13, font_color=C_WHITE)

# Embedder
add_box(sl, Inches(4.5), Inches(4.2), Inches(4.3), Inches(0.55), RGBColor(0x7B, 0x1F, 0xA2),
        "🔢  Embedder (multilingual-e5-large)", font_size=13, font_color=C_WHITE)

# 저장소 — 2개로 분리 (문제 강조)
add_box(sl, Inches(2.0), Inches(5.3), Inches(4.0), Inches(1.5), RGBColor(0xE8, 0x53, 0x3E),
        "NoSQLStore\n(TinyDB — JSON 파일)\n\n문서 메타 + 섹션 내용", font_size=13, font_color=C_WHITE, bold=True)
add_box(sl, Inches(7.3), Inches(5.3), Inches(4.0), Inches(1.5), RGBColor(0xE8, 0x53, 0x3E),
        "VectorStore\n(ChromaDB — 별도 DB)\n\n임베딩 벡터 저장/검색", font_size=13, font_color=C_WHITE, bold=True)

# 한계점 박스
add_box(sl, Inches(0.3), Inches(1.3), Inches(2.0), Inches(5.5), RGBColor(0xFF, 0xEB, 0xEE),
        "", font_size=10, line_color=C_RED)
add_text(sl, Inches(0.35), Inches(1.4), Inches(1.9), Inches(0.35),
         "⚠️ 한계점", font_size=14, font_color=C_RED, bold=True, align=PP_ALIGN.CENTER)
limits = [
    "① DB 2개 분리 운영",
    "  → 동기화 이슈",
    "",
    "② 삭제 시 2곳",
    "   각각 호출 필요",
    "",
    "③ JSON 파일 기반",
    "  → 동시성 취약",
    "",
    "④ 트랜잭션 미지원",
    "  → 데이터 정합성 ↓",
    "",
    "⑤ 확장성 한계",
    "  → 대용량 불리",
]
add_multiline(sl, Inches(0.35), Inches(1.85), Inches(1.85), Inches(5.0), limits,
              font_size=11, font_color=C_RED)


# ═══════════════════════════════════════════════════════════════
# SLIDE 5: TO-BE 아키텍처 (V3)
# ═══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl, C_BG_SLIDE)
slide_title_bar(sl, "03  TO-BE 아키텍처 (V3)", "PostgreSQL 16 + pgvector 통합 설계")
page_number(sl, 5, TOTAL_SLIDES)

# Client
add_box(sl, Inches(4.5), Inches(1.3), Inches(4.3), Inches(0.65), C_PRIMARY,
        "🖥️  Client (Web UI — 동일)", font_size=15, font_color=C_WHITE, bold=True)

# API
add_box(sl, Inches(4.5), Inches(2.3), Inches(4.3), Inches(0.65), RGBColor(0x30, 0x5A, 0x9E),
        "⚡  FastAPI Server (port 8005)", font_size=15, font_color=C_WHITE, bold=True)

# Engine
add_box(sl, Inches(2.5), Inches(3.3), Inches(2.6), Inches(0.55), RGBColor(0x5C, 0x6B, 0xC0),
        "search.py (동일)", font_size=13, font_color=C_WHITE)
add_box(sl, Inches(5.3), Inches(3.3), Inches(2.6), Inches(0.55), RGBColor(0x5C, 0x6B, 0xC0),
        "answer.py (동일)", font_size=13, font_color=C_WHITE)
add_box(sl, Inches(8.1), Inches(3.3), Inches(2.6), Inches(0.55), RGBColor(0x5C, 0x6B, 0xC0),
        "summarizer.py (동일)", font_size=13, font_color=C_WHITE)

# Embedder
add_box(sl, Inches(4.5), Inches(4.2), Inches(4.3), Inches(0.55), RGBColor(0x7B, 0x1F, 0xA2),
        "🔢  Embedder (동일 모델)", font_size=13, font_color=C_WHITE)

# 저장소 — 1개로 통합! (강조)
add_box(sl, Inches(2.5), Inches(5.2), Inches(8.3), Inches(1.8), C_GREEN,
        "PgStore  (PostgreSQL 16 + pgvector 0.8.2)\n하나의 DB에 모든 데이터 통합\n\n"
        "documents 테이블  |  sections 테이블 (+ embedding 컬럼)  |  feedback 테이블\n"
        "HNSW 인덱스  |  ON DELETE CASCADE  |  커넥션 풀링",
        font_size=14, font_color=C_WHITE, bold=True)

# 개선점 박스
add_box(sl, Inches(0.3), Inches(1.3), Inches(2.0), Inches(5.7), RGBColor(0xE8, 0xF5, 0xE9),
        "", font_size=10, line_color=C_GREEN)
add_text(sl, Inches(0.35), Inches(1.4), Inches(1.9), Inches(0.35),
         "✅ 개선 효과", font_size=14, font_color=C_GREEN, bold=True, align=PP_ALIGN.CENTER)
improvements = [
    "① DB 1개로 통합",
    "  → 동기화 불필요",
    "",
    "② CASCADE 삭제",
    "   한 번에 자동 정리",
    "",
    "③ RDBMS 기반",
    "  → 동시성 완벽",
    "",
    "④ ACID 트랜잭션",
    "  → 데이터 정합성 ↑",
    "",
    "⑤ 확장성 우수",
    "  → 대용량 대응 가능",
]
add_multiline(sl, Inches(0.35), Inches(1.85), Inches(1.85), Inches(5.0), improvements,
              font_size=11, font_color=RGBColor(0x2E, 0x7D, 0x32))

# 우측: 핵심 변경 포인트
add_text(sl, Inches(11.2), Inches(1.35), Inches(2), Inches(0.35),
         "🔑 핵심 변경", font_size=13, font_color=C_PRIMARY, bold=True, align=PP_ALIGN.CENTER)
changes = [
    "NoSQLStore",
    "  + VectorStore",
    "     ↓",
    "  PgStore",
    "   (1개 통합)",
]
add_multiline(sl, Inches(11.2), Inches(1.7), Inches(1.8), Inches(2.5), changes,
              font_size=12, font_color=C_DARK)


# ═══════════════════════════════════════════════════════════════
# SLIDE 6: 핵심 구현 — 호출 흐름 (파이프라인)
# ═══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl, C_BG_SLIDE)
slide_title_bar(sl, "04  핵심 구현 — 호출 흐름", "질문 → 임베딩 → 벡터 검색 → LLM 답변 파이프라인")
page_number(sl, 6, TOTAL_SLIDES)

# 파이프라인 흐름도 (좌→우)
steps = [
    ("💬", "사용자 질문",    "자연어 입력",         C_PRIMARY),
    ("🔢", "임베딩 변환",    "E5-large 모델\n1024차원 벡터", RGBColor(0x7B, 0x1F, 0xA2)),
    ("🔍", "벡터 유사도 검색","pgvector HNSW\nCosine 유사도", C_GREEN),
    ("📊", "키워드 매칭",    "테이블명/DR번호\n정확 매칭 부스트", C_ORANGE),
    ("🤖", "LLM 답변 생성",  "GPT-4o-mini\n출처 포함 답변",  RGBColor(0xE5, 0x39, 0x35)),
]

for i, (icon, title, desc, color) in enumerate(steps):
    x = Inches(0.4 + i * 2.55)
    y = Inches(1.5)
    
    # 아이콘 원
    circle = sl.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.55), y, Inches(1.1), Inches(1.1))
    circle.fill.solid()
    circle.fill.fore_color.rgb = color
    circle.line.fill.background()
    tf = circle.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = icon
    p.font.size = Pt(32)
    p.alignment = PP_ALIGN.CENTER
    
    # 제목
    add_text(sl, x, y + Inches(1.25), Inches(2.2), Inches(0.4),
             title, font_size=14, font_color=color, bold=True, align=PP_ALIGN.CENTER)
    
    # 설명
    add_multiline(sl, x + Inches(0.1), y + Inches(1.65), Inches(2.0), Inches(0.9),
                  desc.split("\n"), font_size=11, font_color=C_DARK_GRAY)
    
    # 화살표 (마지막 제외)
    if i < len(steps) - 1:
        add_text(sl, x + Inches(2.1), y + Inches(0.3), Inches(0.5), Inches(0.5),
                 "→", font_size=28, font_color=C_GRAY, align=PP_ALIGN.CENTER)

# 하단: 적재 흐름도
add_rect(sl, Inches(0.3), Inches(4.0), Inches(12.7), Inches(0.04), C_LIGHT)
add_text(sl, Inches(0.5), Inches(4.2), Inches(5), Inches(0.4),
         "📥 문서 적재 파이프라인", font_size=18, font_color=C_PRIMARY, bold=True)

load_steps = [
    ("📄", "DOCX\n파싱", RGBColor(0x42, 0x6F, 0xB1)),
    ("✂️", "섹션\n분할", RGBColor(0x5C, 0x6B, 0xC0)),
    ("🤖", "LLM\n요약", RGBColor(0x7B, 0x1F, 0xA2)),
    ("🔢", "벡터\n임베딩", RGBColor(0x00, 0x96, 0x88)),
    ("💾", "PostgreSQL\n저장", C_GREEN),
]

for i, (icon, label, color) in enumerate(load_steps):
    x = Inches(0.5 + i * 2.55)
    y = Inches(4.9)
    add_box(sl, x, y, Inches(2.0), Inches(1.2), color,
            f"{icon}\n{label}", font_size=14, font_color=C_WHITE, bold=True)
    if i < len(load_steps) - 1:
        add_text(sl, x + Inches(1.9), y + Inches(0.3), Inches(0.7), Inches(0.5),
                 "▶", font_size=22, font_color=C_GRAY, align=PP_ALIGN.CENTER)

# 하단 부연
add_text(sl, Inches(0.5), Inches(6.4), Inches(12), Inches(0.4),
         "* 모든 흐름은 V2와 동일 — 저장소 레이어(PgStore)만 교체, API 엔드포인트 100% 호환",
         font_size=12, font_color=C_GRAY)


# ═══════════════════════════════════════════════════════════════
# SLIDE 7: Before/After 비교
# ═══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl, C_BG_SLIDE)
slide_title_bar(sl, "06  Before / After 비교", "정량적 · 정성적 개선 효과")
page_number(sl, 7, TOTAL_SLIDES)

# 비교 테이블 헤더
add_box(sl, Inches(0.5), Inches(1.35), Inches(3.5), Inches(0.55), C_DARK_GRAY,
        "비교 항목", font_size=15, font_color=C_WHITE, bold=True)
add_box(sl, Inches(4.1), Inches(1.35), Inches(4.3), Inches(0.55), C_RED,
        "V2 (AS-IS)", font_size=15, font_color=C_WHITE, bold=True)
add_box(sl, Inches(8.5), Inches(1.35), Inches(4.3), Inches(0.55), C_GREEN,
        "V3 (TO-BE)", font_size=15, font_color=C_WHITE, bold=True)

rows = [
    ("저장소 구성",      "TinyDB + ChromaDB (2개)",      "PostgreSQL 1개 통합"),
    ("벡터 검색",        "ChromaDB (별도 프로세스)",       "pgvector (DB 내장, HNSW)"),
    ("데이터 정합성",    "파일 기반 — ACID 미지원",        "PostgreSQL ACID 트랜잭션"),
    ("문서 삭제",        "NoSQL삭제 + Vector삭제 (2단계)", "DELETE → CASCADE 자동"),
    ("동시성 처리",      "JSON 파일 Lock 이슈",           "커넥션 풀 + MVCC"),
    ("운영 복잡도",      "DB 2개 백업/관리",              "단일 DB 백업/관리"),
    ("확장성",          "파일 I/O 한계",                 "인덱스 + 파티셔닝 가능"),
    ("API 호환",        "—",                            "100% 동일 엔드포인트"),
]

for i, (item, before, after) in enumerate(rows):
    y = Inches(2.0 + i * 0.62)
    bg = C_WHITE if i % 2 == 0 else RGBColor(0xF0, 0xF4, 0xF8)
    add_box(sl, Inches(0.5), y, Inches(3.5), Inches(0.55), bg,
            item, font_size=13, font_color=C_DARK, bold=True, align=PP_ALIGN.LEFT)
    add_box(sl, Inches(4.1), y, Inches(4.3), Inches(0.55), bg,
            before, font_size=12, font_color=C_RED, align=PP_ALIGN.LEFT)
    add_box(sl, Inches(8.5), y, Inches(4.3), Inches(0.55), bg,
            after, font_size=12, font_color=RGBColor(0x2E, 0x7D, 0x32), bold=True, align=PP_ALIGN.LEFT)

# 하단 요약
add_rect(sl, Inches(0.5), Inches(7.0) - Inches(0.15), Inches(12.3), Inches(0.05), C_PRIMARY)


# ═══════════════════════════════════════════════════════════════
# SLIDE 8: 향후 계획
# ═══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl, C_BG_SLIDE)
slide_title_bar(sl, "08  향후 계획", "확장 로드맵")
page_number(sl, 8, TOTAL_SLIDES)

# 로드맵 타임라인
phases = [
    ("Phase 1 ✅", "현재 완료", [
        "PostgreSQL 16 + pgvector 통합",
        "문서 적재/검색/Q&A 정상 동작",
        "V2 → V3 100% 기능 호환",
    ], C_GREEN),
    ("Phase 2", "단기 (1~2개월)", [
        "전체 문서 일괄 마이그레이션",
        "검색 정확도 튜닝 (가중치 조정)",
        "사용자 피드백 기반 답변 개선",
    ], C_PRIMARY),
    ("Phase 3", "중기 (3~6개월)", [
        "팀 내 공유 및 확산",
        "Oracle DB 연동 조회 기능 강화",
        "변경이력 추적 (문서 버전 비교)",
    ], C_ORANGE),
    ("Phase 4", "장기", [
        "타 시스템 설계문서 확장 적용",
        "Slack/Teams 봇 연동",
        "자동 변경영향도 분석",
    ], RGBColor(0x7B, 0x1F, 0xA2)),
]

for i, (phase, period, items, color) in enumerate(phases):
    x = Inches(0.4 + i * 3.2)
    # 상단 Phase 뱃지
    add_box(sl, x, Inches(1.5), Inches(2.9), Inches(0.55), color,
            phase, font_size=16, font_color=C_WHITE, bold=True)
    add_text(sl, x, Inches(2.15), Inches(2.9), Inches(0.35),
             period, font_size=12, font_color=C_GRAY, align=PP_ALIGN.CENTER)
    # 항목들
    card = add_rect(sl, x, Inches(2.6), Inches(2.9), Inches(3.0), C_WHITE, C_LIGHT, Pt(1))
    bullet_lines = [f"•  {item}" for item in items]
    add_multiline(sl, x + Inches(0.15), Inches(2.75), Inches(2.6), Inches(2.8),
                  bullet_lines, font_size=12, font_color=C_DARK)

# 하단 메시지
add_text(sl, Inches(0.5), Inches(6.2), Inches(12), Inches(0.5),
         "\"설계문서를 읽는 시간을 줄이고, 개발에 집중할 수 있는 환경을 만들겠습니다.\"",
         font_size=18, font_color=C_PRIMARY, bold=True, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════
# 저장
# ═══════════════════════════════════════════════════════════════
out_path = r"C:\Users\my\ia-chatbot-v3\IA_챗봇_V3_발표자료.pptx"
prs.save(out_path)
print(f"✅ PPT 생성 완료: {out_path}")
print(f"   슬라이드: {TOTAL_SLIDES}장")
